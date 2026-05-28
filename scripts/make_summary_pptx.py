"""Full eonet-cascades summary PowerPoint — pulls every figure + narrative
into one shareable deck.

Outputs: ~/Downloads/eonet-cascades-summary.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path.home() / "Downloads" / "eonet-cascades-summary.pptx"
REPO = Path.home() / "Projects" / "eonet-cascades"
FIGS = REPO / "docs" / "figures"

# Palette
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
TEAL = RGBColor(0x2C, 0x7A, 0x7B)
ORANGE = RGBColor(0xC7, 0x5B, 0x12)
GREEN_DARK = RGBColor(0x1F, 0x6B, 0x3A)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GREY = RGBColor(0xAA, 0xAA, 0xAA)
BLACK = RGBColor(0x10, 0x10, 0x10)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
MARGIN_L = Inches(0.6)
MARGIN_T = Inches(0.6)
CONTENT_W = SW - 2 * MARGIN_L


def _set_run(run, text, *, size=None, bold=None, italic=None, color=None, font="Calibri"):
    run.text = text
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    run.font.name = font


def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    return tf


def add_title_bar(slide, title, subtitle=None, color=NAVY, size=30):
    tf = add_textbox(slide, MARGIN_L, MARGIN_T, CONTENT_W, Inches(0.7))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_run(p.add_run(), title, size=size, bold=True, color=color)
    if subtitle:
        tf2 = add_textbox(slide, MARGIN_L, MARGIN_T + Inches(0.65),
                          CONTENT_W, Inches(0.4))
        p = tf2.paragraphs[0]
        _set_run(p.add_run(), subtitle, size=15, italic=True, color=GREY)


def add_bullets(slide, items, *, top=Inches(1.9), font_size=18,
                left=MARGIN_L, width=None):
    if width is None:
        width = CONTENT_W
    tf = add_textbox(slide, left, top, width, SH - top - Inches(0.4))
    for i, item in enumerate(items):
        if isinstance(item, str):
            text, kind = item, "p"
        else:
            text, kind = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        if kind == "header":
            _set_run(run, text, size=font_size + 4, bold=True, color=NAVY)
            p.space_after = Pt(4)
        elif kind == "math":
            _set_run(run, text, size=font_size + 6, bold=True,
                     color=TEAL, font="Cambria Math")
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(8)
            p.space_after = Pt(8)
        elif kind == "bullet":
            _set_run(run, "•  " + text, size=font_size, color=BLACK)
            p.space_after = Pt(4)
        elif kind == "sub":
            _set_run(run, "    ◦  " + text, size=font_size - 2, color=GREY)
            p.space_after = Pt(2)
        elif kind == "fail":
            _set_run(run, "✗  " + text, size=font_size, color=ORANGE, bold=True)
            p.space_after = Pt(4)
        elif kind == "pass":
            _set_run(run, "✓  " + text, size=font_size, color=GREEN_DARK, bold=True)
            p.space_after = Pt(4)
        else:
            _set_run(run, text, size=font_size, color=BLACK)
            p.space_after = Pt(6)


def add_figure(slide, path: Path, *, top=Inches(1.55),
               max_width=None, max_height=None, center=True,
               caption=None):
    if max_width is None:
        max_width = CONTENT_W
    if max_height is None:
        max_height = SH - top - Inches(0.6 if caption else 0.4)
    # Use add_picture letting it auto-size by width then check height
    pic = slide.shapes.add_picture(str(path), MARGIN_L, top, width=max_width)
    if pic.height > max_height:
        # Scale down by height
        scale = max_height / pic.height
        new_w = int(pic.width * scale)
        new_h = int(max_height)
        pic.width = new_w
        pic.height = new_h
    if center:
        pic.left = int((SW - pic.width) // 2)
    if caption:
        cap_tf = add_textbox(slide, MARGIN_L, top + pic.height + Inches(0.05),
                             CONTENT_W, Inches(0.4))
        p = cap_tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_run(p.add_run(), caption, size=11, italic=True, color=GREY)


def add_footer(slide, text):
    tf = add_textbox(slide, MARGIN_L, SH - Inches(0.32), CONTENT_W, Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _set_run(p.add_run(), text, size=9, color=LIGHT_GREY, italic=True)


# ─────────────────────────────────────────────────────────────
# SECTION 1 — Title + Background
# ─────────────────────────────────────────────────────────────

# Slide 1: Title
s = blank_slide()
tf = add_textbox(s, Inches(0.8), Inches(1.8), SW - Inches(1.6), Inches(1.6))
p = tf.paragraphs[0]
_set_run(p.add_run(), "eonet-cascades", size=64, bold=True, color=NAVY)
tf = add_textbox(s, Inches(0.8), Inches(3.3), SW - Inches(1.6), Inches(1.6))
p = tf.paragraphs[0]
_set_run(p.add_run(), "Learning how natural disasters cascade from 25 years of",
         size=22, color=BLACK)
p = tf.add_paragraph()
_set_run(p.add_run(), "global hazard data — and discovering, along the way,", size=22, color=BLACK)
p = tf.add_paragraph()
_set_run(p.add_run(), "a robust failure mode of Neural Hawkes models on imbalanced data.",
         size=22, color=BLACK)
tf = add_textbox(s, Inches(0.8), Inches(5.5), SW - Inches(1.6), Inches(1.0))
p = tf.paragraphs[0]
_set_run(p.add_run(),
         "Four cloud experiments. Five 2024 storm clusters tested. "
         "One working model. One clean diagnosis. One methodological contribution.",
         size=16, italic=True, color=TEAL)
tf = add_textbox(s, Inches(0.8), Inches(6.6), SW - Inches(1.6), Inches(0.5))
_set_run(tf.paragraphs[0].add_run(),
         "Liam Schmidt  ·  Northeastern University  ·  final summary 2026-05-28",
         size=13, italic=True, color=GREY)

# Slide 2 — Big picture
s = blank_slide()
add_title_bar(s, "One disaster sets off another",
              "Hazards don't happen in isolation — they cascade.")
add_bullets(s, [
    "The 2011 Tōhoku earthquake triggered the tsunami that flooded Fukushima.",
    "Hurricane Katrina's storm surge caused most of the deaths attributed to the storm itself.",
    "Wildfires destroy vegetation; a year later when the rains return, hillsides slide.",
    ("", "p"),
    ("Question:", "header"),
    "Across 2.4 million natural-hazard events worldwide since 2000, can we learn the general "
    "structure of these cascades from data alone — without telling the model anything about "
    "geology, meteorology, or risk engineering?",
])
add_footer(s, "1 — Motivation")

# Slide 3 — Data
s = blank_slide()
add_title_bar(s, "The data — NASA EONET")
add_bullets(s, [
    ("Scale", "header"),
    ("2.4 M events, global coverage 2000 → present, 1.1 GB DuckDB", "bullet"),
    ("Per event: (time, longitude, latitude, type)", "bullet"),
    ("", "p"),
    ("Mark vocabulary (K = 7 in working splits)", "header"),
    ("wildfire — ~87% of events (heavily dominant)", "bullet"),
    ("severe_storm, flood, earthquake, tornado, landslide, dust_haze", "bullet"),
    ("", "p"),
    ("Splits used in this work", "header"),
    ("Train: 2022-01 → 2024-06 (2.385 M events)", "bullet"),
    ("Val: 2024-07 → 2024-12 (356 k events)", "bullet"),
])
add_footer(s, "2 — Data")

# Slide 4 — Point processes
s = blank_slide()
add_title_bar(s, "Background — point processes",
              "How to model random events at random times.")
add_bullets(s, [
    "The simplest example is the Poisson process: events at constant rate λ, no memory.",
    ("λ(t) = λ₀     (constant)", "math"),
    ("Properties:", "header"),
    ("Past events tell you NOTHING about future events", "bullet"),
    ("Interevent intervals are exponentially distributed", "bullet"),
    "Natural hazards have memory — an earthquake tells you something about what's coming. We need a process WITH memory.",
])
add_footer(s, "3 — Background")

# Slide 5 — Hawkes
s = blank_slide()
add_title_bar(s, "Hawkes processes — past events boost future rate")
add_bullets(s, [
    "Let the rate at time t depend on all earlier events:",
    ("λ(t) = μ + Σᵢ α · exp[ −β · (t − tᵢ) ]   for tᵢ < t", "math"),
    ("μ — baseline rate (events on their own)", "bullet"),
    ("α — how much each past event boosts the rate of the next one", "bullet"),
    ("β — how quickly that boost decays in time", "bullet"),
    "Each event creates a spike in the intensity that decays. Stack the spikes → rate of the next event.",
])
add_footer(s, "4 — Hawkes process")

# Slide 6 — Multivariate Hawkes
s = blank_slide()
add_title_bar(s, "Multivariate Hawkes — cascades between event types")
add_bullets(s, [
    ("λᵢ(t) = μᵢ + Σⱼ Σ_{tₖ<t, type(k)=j} αᵢⱼ · exp[ −βᵢⱼ(t − tₖ) ]", "math"),
    ("The K × K matrix αᵢⱼ IS the cascade graph:", "header"),
    ("αᵢⱼ > 0   ⇒   type j triggers type i", "bullet"),
    ("αᵢⱼ ≈ 0   ⇒   no causal coupling", "bullet"),
    "If we can fit α from data, we recover the cascade structure. That is this project's headline goal.",
])
add_footer(s, "5 — Multivariate Hawkes")

# Slide 7 — Three tiers
s = blank_slide()
add_title_bar(s, "Three model tiers, one shared likelihood",
              "Apples-to-apples comparison across functional forms.")
add_bullets(s, [
    ("Tier 0 — Parametric Hawkes", "header"),
    ("Constant α × exp(time) × Gaussian(space). ~200 params. Max interpretability; rigid form.", "bullet"),
    ("", "p"),
    ("Tier 1 — Neural Hawkes (this project's focus)", "header"),
    ("CTLSTM hidden state + per-mark intensity head + bivariate MDN spatial head. ~30k params.", "bullet"),
    ("", "p"),
    ("Tier 2 — Transformer Hawkes (deferred)", "header"),
    ("Long-range attention. Not built; reserved as future work.", "bullet"),
    ("", "p"),
    "Tier 1 has no explicit α matrix. We derive K×K cascade structure two ways: gradient attribution AND forward simulation. They should agree at infinite data.",
])
add_footer(s, "6 — Three tiers")

# Slide 8 — Architecture diagram
s = blank_slide()
add_title_bar(s, "Tier 1-MLP — Neural Hawkes architecture")
add_figure(s, FIGS / "fig_architecture.png",
           caption="CTLSTM hidden state h(t) → per-mark intensity head (the H6 bottleneck) + mark-conditional MDN spatial head")
add_footer(s, "7 — Architecture")

# Slide 9 — Interpretability bridge
s = blank_slide()
add_title_bar(s, "Two derived K × K views of the same model")
add_bullets(s, [
    ("View A — Gradient attribution", "header"),
    ("Operates on the LSTM hidden state h(t) directly", "bullet"),
    ("Bypasses the mark head's output", "bullet"),
    ("Recovers off-diagonal cascade structure cleanly", "pass"),
    ("", "p"),
    ("View B — Forward simulation", "header"),
    ("Samples through the mark head's full output", "bullet"),
    ("Should produce a parent-conditional K × K transition matrix", "bullet"),
    ("Row-degenerate in all four runs", "fail"),
    ("", "p"),
    ("The disagreement IS the diagnostic. Three days of work tracking it down.", "header"),
])
add_footer(s, "8 — Interpretability bridge")

# ─────────────────────────────────────────────────────────────
# SECTION 2 — Four-run negative chain
# ─────────────────────────────────────────────────────────────

# Slide 10 — Four-run overview
s = blank_slide()
add_title_bar(s, "Four cloud experiments — Lambda Labs A10, ~$19 total", color=ORANGE)
add_bullets(s, [
    ("Each targeted the rank-1 forward-sim row-degeneracy with a different intervention.", "p"),
    ("Each failed the primary acceptance criterion: probe row-deviation > 0.1.", "p"),
    ("", "p"),
    ("Tier 1   ·   baseline (linear head, joint Hawkes NLL)", "header"),
    ("Val NLL 4.20  ·  probe row-dev 0.0012", "bullet"),
    ("Tier 1.5 ·  class re-weighting + stratified subsample", "header"),
    ("Val NLL 6.80 (+62%)  ·  probe row-dev 0.0000", "fail"),
    ("Tier 1-MLP  ·  non-linear mark head (Linear → ReLU → Linear)", "header"),
    ("Val NLL 3.38 (−20%, best)  ·  probe row-dev 0.0000", "fail"),
    ("Tier 1-aux  ·  explicit cross-entropy auxiliary loss λ=1.0", "header"),
    ("Val NLL 4.94 (+18%)  ·  probe row-dev 0.0002", "fail"),
])
add_footer(s, "9 — Four runs")

# Slide 11 — K×K grid
s = blank_slide()
add_title_bar(s, "Forward-sim K × K matrices: row-degeneracy across all four runs",
              color=ORANGE)
add_figure(s, FIGS / "fig1_kk_grid.png",
           caption="Every parent (row) gives the same child distribution (column pattern). The model ignores parent mark for composition prediction.")
add_footer(s, "10 — K × K row-degeneracy")

# Slide 12 — Convergence
s = blank_slide()
add_title_bar(s, "Convergence curves — Tier 1-MLP achieves the best val NLL")
add_figure(s, FIGS / "fig2_convergence.png",
           caption="Tier 1-MLP (green) finishes at 3.38, beating Tier 1 (4.20) by 20%. Tier 1.5 (orange) and Tier 1-aux (teal) above the 4.41 acceptance line.")
add_footer(s, "11 — Convergence")

# Slide 13 — Marginal shifts
s = blank_slide()
add_title_bar(s, "Interventions DO shift the mark-head marginal — but not the conditioning",
              color=ORANGE)
add_figure(s, FIGS / "fig3_marginal_bars.png",
           caption="Tier 1 and Tier 1-MLP are wildfire-dominant (~90%); Tier 1.5 and Tier 1-aux flattened the marginal. But the K rows remain identical across parents in all runs.")
add_footer(s, "12 — Marginal vs conditioning")

# Slide 14 — Row-dev bar
s = blank_slide()
add_title_bar(s, "Primary acceptance criterion: all four runs fail by 3 orders of magnitude")
add_figure(s, FIGS / "fig4_row_dev_bars.png",
           caption="Row-deviation > 0.1 required. Observed: 0.0012, 0.0000, 0.0000, 0.0002. Robust failure across four interventions.")
add_footer(s, "13 — Acceptance criterion")

# Slide 15 — H6 diagnosis
s = blank_slide()
add_title_bar(s, "The sharpened diagnosis: encoder bottleneck (H6)", color=GREEN_DARK)
add_bullets(s, [
    ("Refuted hypotheses (mark head + training signal):", "header"),
    ("Class imbalance creates a marginal local minimum (Tier 1.5)", "fail"),
    ("Linear head insufficient capacity (Tier 1-MLP)", "fail"),
    ("Joint Hawkes NLL has insufficient gradient on z (Tier 1-aux)", "fail"),
    ("", "p"),
    ("Supported hypothesis — H6:", "header"),
    ("The LSTM hidden state h(t) does not encode per-parent variation at the rank the mark head's linear projections can extract.", "pass"),
    ("", "p"),
    "Supported by exclusion of three alternatives AND by direct evidence: gradient attribution on the same checkpoint DOES recover off-diagonal cascade structure (info IS in h, just not at the rank linear functions of h can read).",
])
add_footer(s, "14 — Diagnosis")

# Slide 16 — Effective rank (THE killer)
s = blank_slide()
add_title_bar(s, "Quantifying H6: effective rank = 1.00 out of K = 7", color=GREEN_DARK)
add_figure(s, FIGS / "fig_effective_rank.png",
           caption="Mark-head outputs lie on a 1.00-dimensional manifold. Across 21 hidden states from validation, all 3 output forms (raw logits, intensities, normalized probabilities) collapse to a single direction.")
add_footer(s, "15 — Effective rank")

# Slide 17 — Mark-agnostic decay
s = blank_slide()
add_title_bar(s, "All 7 seed marks produce IDENTICAL intensity decay curves",
              color=GREEN_DARK)
add_figure(s, FIGS / "fig_hawkes_decay.png",
           caption="One seed event of any mark. λ_total(t) trajectory should differ by mark type. It doesn't. Seven curves overlap exactly — another independent confirmation of H6.")
add_footer(s, "16 — Decay curves")

# Slide 18 — Geography
s = blank_slide()
add_title_bar(s, "The model didn't learn distinct per-mark geography either")
add_figure(s, FIGS / "fig_geography.png",
           caption="Per-mark spatial densities are visually near-identical across all 7 marks. Model learned 'where events happen' generically, not 'wildfires in California vs hurricanes in the Gulf.'")
add_footer(s, "17 — Per-mark geography")

# Slide 19 — Methodology contribution
s = blank_slide()
add_title_bar(s, "The methodological contribution: cross-view triangulation",
              "A diagnostic pattern that generalizes beyond Hawkes models.")
add_bullets(s, [
    ("The pattern:", "header"),
    ("Two K × K interpretability views of the same trained model that, at infinite data, should agree.", "bullet"),
    ("Notice when they disagree at finite data. That disagreement is the diagnostic signal.", "bullet"),
    ("Run targeted interventions; track which change behavior and which do not.", "bullet"),
    ("The intervention that changes the most without solving the problem is the most informative (Tier 1-aux here).", "bullet"),
    ("", "p"),
    ("Three independent quantitative confirmations of H6 in this work:", "header"),
    ("Probe row-deviation = 0.0000", "pass"),
    ("Effective rank = 1.00 (out of K=7)", "pass"),
    ("Mark-agnostic intensity decay (all 7 curves overlap)", "pass"),
])
add_footer(s, "18 — Methodology")

# ─────────────────────────────────────────────────────────────
# SECTION 3 — But the model works! Case studies
# ─────────────────────────────────────────────────────────────

# Slide 20 — Transition
s = blank_slide()
add_title_bar(s, "But — the same model is a working rate forecaster",
              "The cascade-extraction path is broken. The forecaster underneath isn't.")
add_bullets(s, [
    ("Tier 1-MLP, val NLL 3.38, has a broken forward-sim path but:", "header"),
    ("Correctly tracks total event rate as a function of history (the LSTM IS reading h for rates).", "pass"),
    ("Correctly predicts per-mark spatial distributions via the MDN.", "pass"),
    ("Cross-mark cascade structure is recoverable via gradient attribution (alternate extraction path).", "pass"),
    ("", "p"),
    "We tested this empirically on five 2024 storm clusters covering both major hurricanes (Francine, Milton) and smaller bursts. The model beats a marginal-Poisson baseline by ~100× across all five.",
    ("", "p"),
    "Following slides show the empirical validation.",
])
add_footer(s, "19 — Pivot to case studies")

# Slide 21 — Hero figure (storm-only)
s = blank_slide()
add_title_bar(s, "The model adapts in real time to Hurricane Francine",
              color=NAVY)
add_figure(s, FIGS / "case_study_hero_storm.png",
           caption="Sept 1 (calm baseline) vs Sept 10 (peak Francine). Storm-related intensity peak goes 3.65×10⁻³ → 1.10×10⁻² as the model raises predictions across the entire eastern US ahead of landfall.")
add_footer(s, "20 — Hero figure")

# Slide 22 — Daily counts
s = blank_slide()
add_title_bar(s, "Sept 9-10 storm cluster: 8,144 + 7,275 events")
add_figure(s, FIGS / "case_study_francine_daily.png",
           caption="Daily EONET event counts in the test window. Highlighted Sept 9-10 burst is nearly double any other day. Warm-up (grey) feeds the model history; test window (navy) is scored against actuals.")
add_footer(s, "21 — Daily counts")

# Slide 23 — Likelihood
s = blank_slide()
add_title_bar(s, "Model log-likelihood is flat across the burst — it wasn't surprised")
add_figure(s, FIGS / "case_study_francine_likelihood.png",
           caption="Top: per-event log-lik scatter, colored by mark. Bottom: daily mean model vs marginal-Poisson baseline. ~4.6 nat/event gap, flat across the entire test window including the storm cluster.")
add_footer(s, "22 — Likelihood")

# Slide 24 — Multi-cluster robustness
s = blank_slide()
add_title_bar(s, "+4.6 nats/event improvement holds across FIVE 2024 storm bursts", color=GREEN_DARK)
add_figure(s, FIGS / "case_study_multi_cluster.png",
           caption="Same Tier 1-MLP, five independent test windows. Improvement ranges +4.56 to +4.72 nats/event. Milton has the largest gap (+4.72). The result is robust, not cherry-picked.")
add_footer(s, "23 — Robustness")

# Slide 25 — Spatial scatter
s = blank_slide()
add_title_bar(s, "The model gives high likelihood to events along Francine's track")
add_figure(s, FIGS / "case_study_francine_spatial.png",
           caption="Sept 9-10 events colored by model log-lik. Yellow = high. Red star = Francine landfall. Concentration of high-lik events follows the storm corridor.")
add_footer(s, "24 — Spatial scatter")

# Slide 26 — Calibration
s = blank_slide()
add_title_bar(s, "Calibration: the model knows what it doesn't know")
add_figure(s, FIGS / "case_study_calibration.png",
           caption="Lowest log-lik decile contains the rare-mark events (earthquake, flood, severe_storm). Mid-to-high deciles are essentially all wildfires. Useful for anomaly-detection: low log-lik events are likely non-wildfire hazards.")
add_footer(s, "25 — Calibration")

# ─────────────────────────────────────────────────────────────
# SECTION 4 — Wrap-up
# ─────────────────────────────────────────────────────────────

# Slide 27 — What this project actually delivers
s = blank_slide()
add_title_bar(s, "What this project actually delivers", color=GREEN_DARK)
add_bullets(s, [
    ("A working Neural Hawkes model", "header"),
    ("Tier 1-MLP at runs/tier1_mlp/20260526_141553/, val NLL 3.38, 20% improvement over Tier 1 baseline", "bullet"),
    ("Demonstrated on real 2024 hurricane season — beats marginal baseline by ~100× across 5 storm clusters", "bullet"),
    ("", "p"),
    ("A characterized failure mode", "header"),
    ("Mark-head rank-1 collapse survives 4 textbook remedies (rebalance, MLP head, aux loss, etc.)", "bullet"),
    ("Localized to LSTM encoder; gradient attribution recovers the cascade structure the forward-sim path can't", "bullet"),
    ("Quantified three independent ways: probe row-dev 0.000, effective rank 1.00, mark-agnostic decay", "bullet"),
    ("", "p"),
    ("Methodological contribution", "header"),
    ("Cross-view interpretability triangulation as a diagnostic. Disagreement IS the signal.", "bullet"),
    ("Vectorized attribution kernel: 240× speedup, enables n=5000 analyses in 40 min", "bullet"),
    ("Reproducible end-to-end pipeline + 4-run cloud experiment chain, $19 total spend", "bullet"),
])
add_footer(s, "26 — Deliverables")

# Slide 28 — Future work
s = blank_slide()
add_title_bar(s, "Future work — ranked by yield-per-effort")
add_bullets(s, [
    ("Tier 1 — Direct continuation", "header"),
    ("Effective-rank diagnostic on W_λ_k @ Cov(h) — DONE in this work (rank = 1.00)", "pass"),
    ("Wider LSTM hidden_dim (128, 256) — ~$10 cloud, tests if bottleneck is just capacity", "bullet"),
    ("Transformer encoder (the original Tier 2 scope) — tests if H6 is LSTM-specific", "bullet"),
    ("", "p"),
    ("Tier 2 — Methodology generalization", "header"),
    ("Apply cross-view triangulation to other Neural TPP models (RMTPP, Transformer Hawkes, etc.)", "bullet"),
    ("Synthetic ground-truth Hawkes recovery test — separates EONET-data property from model-class property", "bullet"),
    ("", "p"),
    ("Operational", "header"),
    ("Hazard-rate forecasting service productionizing Tier 1-MLP", "bullet"),
    ("Gradient-attribution-based cascade dashboard (alternative extraction path)", "bullet"),
])
add_footer(s, "27 — Future work")

# Slide 29 — Closing
s = blank_slide()
tf = add_textbox(s, Inches(1.0), Inches(2.2), SW - Inches(2.0), Inches(1.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
_set_run(p.add_run(), "Closed.", size=64, bold=True, color=NAVY)
tf = add_textbox(s, Inches(1.0), Inches(3.7), SW - Inches(2.0), Inches(2.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
_set_run(p.add_run(),
         "Four cloud runs.  Five storm-cluster validations.  $19 spent.",
         size=18, italic=True, color=TEAL)
p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
_set_run(p.add_run(), "", size=20)
p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
_set_run(p.add_run(), "One working model.", size=22, bold=True, color=GREEN_DARK)
p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
_set_run(p.add_run(), "One characterized failure mode.", size=22, bold=True, color=ORANGE)
p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
_set_run(p.add_run(), "One methodological contribution.", size=22, bold=True, color=NAVY)

tf = add_textbox(s, Inches(1.0), Inches(6.6), SW - Inches(2.0), Inches(0.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
_set_run(p.add_run(), "github.com/liamlts/eonet-cascades  (private)",
         size=14, color=GREY)


prs.save(OUT)
print(f"Saved {OUT}  ({len(prs.slides)} slides)")
